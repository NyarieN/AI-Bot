import os
import json
import re
import math
import tempfile
import numpy as np
import joblib
from flask import Flask, render_template, request, jsonify, redirect, flash
from werkzeug.utils import secure_filename
from dotenv import load_dotenv
from cryptography.fernet import Fernet
from io import BytesIO, StringIO
from sklearn.metrics.pairwise import cosine_similarity

# --- Base Directories ---
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
UPLOAD_FOLDER = os.path.join(BASE_DIR, "uploads")
EMBEDDINGS_FOLDER = os.path.join(BASE_DIR, "embeddings")

# Create base folders if they don't exist
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(EMBEDDINGS_FOLDER, exist_ok=True)

# Optional libs (used only if available)
try:
    import PyPDF2
except Exception:
    PyPDF2 = None

try:
    import openpyxl
except Exception:
    openpyxl = None

try:
    from docx import Document
except Exception:
    Document = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    # TF-IDF fallback (local)
    from sklearn.feature_extraction.text import TfidfVectorizer
except Exception:
    TfidfVectorizer = None

# Optional: numpy for faster vector math (not required)
try:
    import numpy as np
except Exception:
    np = None

# OpenAI (optional)
try:
    from openai import OpenAI
except Exception:
    OpenAI = None

load_dotenv()

# -----------------------
# Config
# -----------------------
app = Flask(__name__)
app.secret_key = os.getenv("FLASK_SECRET_KEY", "anysecret")

ROOT_UPLOADS = "uploads"
ROOT_EMBEDDINGS = "embeddings"
os.makedirs(ROOT_UPLOADS, exist_ok=True)
os.makedirs(ROOT_EMBEDDINGS, exist_ok=True)

ALLOWED_EXTENSIONS = {"pdf", "txt", "docx", "doc", "xlsx", "xls", "csv", "pptx", "ppt"}

# Encryption key setup
ENCRYPTION_KEY = os.getenv("ENCRYPTION_KEY")
if not ENCRYPTION_KEY:
    # generate but do NOT print in production
    ENCRYPTION_KEY = Fernet.generate_key().decode()
fernet = Fernet(ENCRYPTION_KEY.encode() if isinstance(ENCRYPTION_KEY, str) else ENCRYPTION_KEY)

# OpenAI client available only if key and OpenAI package available
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
USE_OPENAI = bool(OPENAI_API_KEY and OpenAI is not None)
client = OpenAI(api_key=OPENAI_API_KEY) if USE_OPENAI else None

# Models
EMBEDDING_MODEL = "text-embedding-3-small"
CHAT_MODEL = "gpt-4o-mini"

# -----------------------
# Helpers
# -----------------------
def allowed_file(filename):
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS

def encrypt_file(file_path):
    with open(file_path, "rb") as f:
        data = f.read()
    enc = fernet.encrypt(data)
    with open(file_path, "wb") as f:
        f.write(enc)

def decrypt_file(file_path):
    with open(file_path, "rb") as f:
        data = f.read()
    return fernet.decrypt(data)

def safe_business_folder(base_root, business_id):
    """Return path for the business folder under base_root and ensure it's created."""
    safe_id = secure_filename(business_id or "default")
    path = os.path.join(base_root, safe_id)
    os.makedirs(path, exist_ok=True)
    return path

def read_file_bytes(path):
    with open(path, "rb") as f:
        return f.read()

def extract_text_from_file_on_disk(file_path):
    """
    Extract text from a raw (unencrypted) file on disk.
    This is used immediately after saving a newly uploaded file (before encryption).
    """
    ext = file_path.rsplit(".", 1)[1].lower()
    try:
        data = read_file_bytes(file_path)
        if ext == "pdf" and PyPDF2:
            f = BytesIO(data)
            reader = PyPDF2.PdfReader(f)
            return "\n".join([page.extract_text() or "" for page in reader.pages])
        elif ext == "txt":
            return data.decode(errors="ignore")
        elif ext in ("docx", "doc") and Document:
            doc = Document(BytesIO(data))
            return "\n".join([p.text for p in doc.paragraphs])
        elif ext in ("xlsx", "xls") and openpyxl:
            wb = openpyxl.load_workbook(BytesIO(data), data_only=True)
            text = ""
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                for row in ws.iter_rows(values_only=True):
                    text += " ".join([str(cell) for cell in row if cell]) + "\n"
            return text
        elif ext == "csv":
            f = StringIO(data.decode(errors="ignore"))
            import csv as _csv
            reader = _csv.reader(f)
            return "\n".join([" ".join(row) for row in reader])
        elif ext in ("pptx", "ppt") and Presentation:
            prs = Presentation(BytesIO(data))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        else:
            # Unknown type or missing optional libs — try decode
            return data.decode(errors="ignore")
    except Exception as e:
        app.logger.warning(f"Failed to extract text from {file_path}: {e}")
        return ""

def split_text_into_chunks(text, chunk_size=500, overlap=50):
    """
    Split text into chunks of roughly chunk_size characters with optional overlap.
    This is simple char-based chunking (sufficient for embeddings/TF-IDF).
    """
    if not text:
        return []
    text = re.sub(r"\s+", " ", text).strip()
    chunks = []
    start = 0
    text_len = len(text)
    while start < text_len:
        end = start + chunk_size
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        start = end - overlap
        if start < 0:
            start = 0
    return chunks

def dot(a, b):
    return sum(x * y for x, y in zip(a, b))

def norm(a):
    return math.sqrt(sum(x * x for x in a))

def cosine_similarity(a, b):
    # a and b are lists of floats
    try:
        denom = norm(a) * norm(b)
        if denom == 0:
            return 0.0
        return dot(a, b) / denom
    except Exception:
        return 0.0

# -----------------------
# Embedding helpers
# -----------------------
def generate_openai_embeddings_for_chunks(chunks):
    """Return list of dicts: {'text': chunk, 'embedding': [floats]}"""
    if not USE_OPENAI:
        raise RuntimeError("OpenAI not available for embeddings")
    out = []
    # Batch if needed; keep simple per-chunk calls to avoid large payloads
    for chunk in chunks:
        try:
            resp = client.embeddings.create(model=EMBEDDING_MODEL, input=chunk)
            vec = resp.data[0].embedding
            out.append({"text": chunk, "embedding": vec})
        except Exception as e:
            app.logger.warning(f"Embedding failed for a chunk: {e}")
    return out

# -----------------------
# Admin portal
# -----------------------
@app.route("/admin", methods=["GET", "POST"])
def admin_upload():
    # temporary simple auth via query param (existing behavior)
    key = request.args.get("key")
    if key != os.getenv("ADMIN_KEY", "supersecret123"):
        return "Unauthorized", 403

    if request.method == "POST":
        # sanitize business id
        raw_business_id = request.form.get("business_id", "").strip().lower()
        if not raw_business_id:
            flash("Missing business_id", "error")
            return redirect(request.url)
        business_id = secure_filename(raw_business_id)

        files = request.files.getlist("files")
        if not files:
            flash("No files uploaded", "error")
            return redirect(request.url)

        business_upload_folder = safe_business_folder(ROOT_UPLOADS, business_id)
        business_embeddings_folder = safe_business_folder(ROOT_EMBEDDINGS, business_id)

        all_texts = []
        uploaded_files = []
        skipped_files = []

        # Save uploaded files, extract text before encrypting
        for file in files:
            filename = file.filename or ""
            if not filename or not allowed_file(filename):
                skipped_files.append(filename or "unknown")
                continue
            filename = secure_filename(filename)
            path = os.path.join(business_upload_folder, filename)
            file.save(path)
            # Extract text from raw (unencrypted) file
            text = extract_text_from_file_on_disk(path)
            # Encrypt file in place
            try:
                encrypt_file(path)
            except Exception as e:
                app.logger.warning(f"Could not encrypt {path}: {e}")
            if text:
                all_texts.append(text)
            uploaded_files.append(filename)

        # Save combined text to all_text.txt and encrypt
        combined_text = "\n".join(all_texts).strip()
        combined_text_path = os.path.join(business_upload_folder, "all_text.txt")
        if combined_text:
            with open(combined_text_path, "w", encoding="utf-8") as f:
                f.write(combined_text)
            encrypt_file(combined_text_path)

        # Always create and save chunks.json (plaintext then encrypted)
        chunks = []
        for t in all_texts:
            chunks.extend(split_text_into_chunks(t, chunk_size=500, overlap=100))

        chunks_path = os.path.join(business_embeddings_folder, "chunks.json")
        with open(chunks_path, "w", encoding="utf-8") as f:
            json.dump(chunks, f)
        encrypt_file(chunks_path)

        # If OpenAI is available, generate semantic embeddings and save, else skip
        embeddings_path = os.path.join(business_embeddings_folder, "embeddings.json")
        if USE_OPENAI and chunks:
            try:
                sem = generate_openai_embeddings_for_chunks(chunks)
                with open(embeddings_path, "w", encoding="utf-8") as f:
                    json.dump(sem, f)
                encrypt_file(embeddings_path)
            except Exception as e:
                app.logger.warning(f"Failed to generate embeddings for {business_id}: {e}")

        # Feedback to admin
        if uploaded_files:
            flash(f"Uploaded: {', '.join(uploaded_files)}", "success")
        if skipped_files:
            flash(f"Skipped (unsupported): {', '.join(skipped_files)}", "error")
        if not uploaded_files:
            flash("No files were saved.", "error")

        return redirect(request.url)

    return render_template("admin.html")

# -----------------------
# Query (FAQ) route
# -----------------------
# -----------------------
# Improved Query Route (Per Business Isolation + Smart Search)
# -----------------------
@app.route("/query", methods=["POST"])
def query_pdf():
    data = request.json
    business_id = data.get("business_id", "").lower().strip()
    question = data.get("question", "").strip()

    if not business_id or not question:
        return jsonify({"answer": "Missing business_id or question."}), 400

    # Paths for this business only
    text_file = os.path.join(UPLOAD_FOLDER, business_id, "all_text.txt")
    embeddings_folder = os.path.join("embeddings", business_id)

    # Validate that business data exists
    if not os.path.exists(text_file):
        return jsonify({"answer": "No documents uploaded for this business."}), 404

    try:
        # First try to use embeddings if they exist
        vectorizer_path = os.path.join(embeddings_folder, "vectorizer.pkl")
        matrix_path = os.path.join(embeddings_folder, "matrix.pkl")
        docs_path = os.path.join(embeddings_folder, "docs.pkl")
        vectors_path = os.path.join(embeddings_folder, "vectors.pkl")

        # Check if we have TF-IDF or OpenAI mode
        use_openai = os.getenv("OPENAI_API_KEY") and os.path.exists(vectors_path)

        if use_openai:
            # --- OpenAI Embeddings Mode ---
            print(f"Using OpenAI embeddings for {business_id}")
            import numpy as np
            with open(docs_path, "rb") as f:
                docs = joblib.load(f)
            with open(vectors_path, "rb") as f:
                vectors = joblib.load(f)

            # Embed the user query
            q_emb = client.embeddings.create(
                input=question[:8000],
                model="text-embedding-3-small"
            ).data[0].embedding

            sims = cosine_similarity([q_emb], vectors)[0]
            best_idx = int(np.argmax(sims))
            best_match = docs[best_idx]
            return jsonify({"answer": best_match})

        elif os.path.exists(vectorizer_path) and os.path.exists(matrix_path):
            # --- Local TF-IDF Mode ---
            print(f"Using local TF-IDF for {business_id}")
            vectorizer = joblib.load(vectorizer_path)
            matrix = joblib.load(matrix_path)
            docs = joblib.load(docs_path)

            q_vec = vectorizer.transform([question])
            sims = cosine_similarity(q_vec, matrix).flatten()

            if len(sims) == 0 or np.max(sims) < 0.05:
                return jsonify({"answer": "I could not find a relevant answer in the documents."})

            best_idx = int(np.argmax(sims))
            best_match = docs[best_idx]
            return jsonify({"answer": best_match})

        else:
            # --- Fallback: old simple keyword matching ---
            print(f"No embeddings found for {business_id}, using basic match")
            all_text = decrypt_file(text_file).decode()
            sentences = re.split(r'(?<=[.!?]) +', all_text)
            best_sentence = max(
                sentences,
                key=lambda s: sum(word in s.lower() for word in question.split()),
                default="I could not find a relevant answer for you in the documents."
            )
            return jsonify({"answer": best_sentence})

    except Exception as e:
        print(f"Query error for {business_id}: {e}")
        return jsonify({"answer": f"I could not find a relevant answer for you. {e}"})

# -----------------------
# AI (chat) route
# -----------------------
@app.route("/ask", methods=["POST"])
def ask_ai():
    data = request.json or {}
    business_id_raw = data.get("business_id", "").strip().lower()
    user_message = data.get("message", "").strip()
    if not business_id_raw:
        return jsonify({"reply": "Missing business_id."})
    if not user_message:
        return jsonify({"reply": "Please send a message to ask."})

    business_id = secure_filename(business_id_raw)

    # If OpenAI is available, call chat API
    if USE_OPENAI:
        try:
            response = client.chat.completions.create(
                model=CHAT_MODEL,
                messages=[
                    {"role": "system", "content": f"You are Kai, assistant for {business_id}."},
                    {"role": "user", "content": user_message}
                ],
            )
            # response structure: response.choices[0].message.content
            reply = response.choices[0].message.content
            return jsonify({"reply": reply})
        except Exception as e:
            app.logger.warning(f"OpenAI chat error: {e}")
            # fallback to local query: run /query logic internally
            # reuse query_pdf logic by constructing a fake request
    # OpenAI not available or failed -> fallback to local FAQ query
    # Reuse query logic: craft a local call to query_pdf
    try:
        # call the query route function directly
        fake_json = {"business_id": business_id, "question": user_message}
        # query_pdf expects request.json; but we can call the logic by importing function or duplicating minimal logic:
        # For simplicity, run the query logic locally here (same fallback chain as /query).
        # We'll attempt to use semantic embeddings if present (but OPENAI not available) -> skip to TF-IDF or keyword.
        # Use same code as /query fallback:
        business_upload_folder = os.path.join(ROOT_UPLOADS, business_id)
        business_embeddings_folder = os.path.join(ROOT_EMBEDDINGS, business_id)
        chunks_enc_path = os.path.join(business_embeddings_folder, "chunks.json")
        all_text_path_enc = os.path.join(business_upload_folder, "all_text.txt")

        if not os.path.exists(all_text_path_enc) or not os.path.exists(chunks_enc_path):
            return jsonify({"reply": "No documents uploaded for this business to answer from."})

        chunks_data = json.loads(decrypt_file(chunks_enc_path).decode(errors="ignore"))

        # TF-IDF fallback
        if TfidfVectorizer and chunks_data:
            try:
                vectorizer = TfidfVectorizer().fit(chunks_data)
                chunk_matrix = vectorizer.transform(chunks_data)
                q_vec = vectorizer.transform([user_message])
                sims = (q_vec @ chunk_matrix.T).toarray()[0]
                if len(sims) > 0:
                    best_idx = int(sims.argmax())
                    best_chunk = chunks_data[best_idx]
                    return jsonify({"reply": best_chunk})
            except Exception:
                pass

        # last fallback: keyword on sentences
        all_text = decrypt_file(all_text_path_enc).decode(errors="ignore")
        sentences = re.split(r'(?<=[.!?])\s+', all_text)
        query_words = [w.lower() for w in re.findall(r"\w+", user_message)]
        def score_sentence(s):
            sw = [w.lower() for w in re.findall(r"\w+", s)]
            if not sw:
                return 0
            match_count = sum(1 for qw in query_words if qw in sw)
            return match_count / (len(sw) + 1)
        best_sentence = max(sentences, key=score_sentence, default="I could not find an answer in the documents.")
        return jsonify({"reply": best_sentence})
    except Exception as e:
        app.logger.warning(f"Local fallback (ask) failed: {e}")
        return jsonify({"reply": "Sorry, something went wrong while processing your request.Please contact customersupport@prismdx.com"})

# -----------------------
# Regenerate Embeddings (Admin Only)
# -----------------------
from sklearn.feature_extraction.text import TfidfVectorizer
import joblib

@app.route("/regenerate_embeddings", methods=["POST"])
def regenerate_embeddings():
    # --- Admin key check ---
    key = request.args.get("key")
    if key != os.getenv("ADMIN_KEY", "supersecret123"):
        return jsonify({"error": "Unauthorized"}), 403

    data = request.get_json()
    business_id = data.get("business_id")
    if not business_id:
        return jsonify({"error": "Missing business_id"}), 400

    business_id = business_id.lower()
    business_folder = os.path.join(UPLOAD_FOLDER, business_id)
    embeddings_folder = os.path.join("embeddings", business_id)
    os.makedirs(embeddings_folder, exist_ok=True)

    text_file = os.path.join(business_folder, "all_text.txt")
    if not os.path.exists(text_file):
        return jsonify({"error": f"No uploaded data found for {business_id}"}), 404

    try:
        # Decrypt the file to read the text
        all_text = decrypt_file(text_file).decode()

        # Split into smaller chunks
        docs = [t.strip() for t in all_text.split("\n") if len(t.strip()) > 5]

        if not docs:
            return jsonify({"error": "No readable text found"}), 400

        # Check if OpenAI API key is set
        openai_key = os.getenv("OPENAI_API_KEY")
        use_openai = bool(openai_key)

        if use_openai:
            # --- Use OpenAI embeddings ---
            print(f"Using OpenAI embeddings for {business_id} ...")
            vectors = []
            for doc in docs:
                try:
                    emb = client.embeddings.create(
                        input=doc[:8000],  # limit length
                        model="text-embedding-3-small"
                    ).data[0].embedding
                    vectors.append(emb)
                except Exception as e:
                    print(f"Embedding failed: {e}")

            joblib.dump(docs, os.path.join(embeddings_folder, "docs.pkl"))
            joblib.dump(vectors, os.path.join(embeddings_folder, "vectors.pkl"))
        else:
            # --- Use TF-IDF (offline mode) ---
            print(f"Using local TF-IDF embeddings for {business_id} ...")
            vectorizer = TfidfVectorizer(stop_words="english", max_features=5000)
            matrix = vectorizer.fit_transform(docs)

            joblib.dump(vectorizer, os.path.join(embeddings_folder, "vectorizer.pkl"))
            joblib.dump(matrix, os.path.join(embeddings_folder, "matrix.pkl"))
            joblib.dump(docs, os.path.join(embeddings_folder, "docs.pkl"))

        return jsonify({
            "status": "success",
            "message": f"Embeddings regenerated for {business_id}",
            "method": "openai" if use_openai else "tfidf"
        })

    except Exception as e:
        print(f"Error regenerating embeddings: {e}")
        return jsonify({"error": str(e)}), 500

# -----------------------
# Home
# -----------------------
@app.route("/")
def home():
    return render_template("index.html")

# -----------------------
# Run
# -----------------------
if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5000))
    app.run(host="0.0.0.0", port=port)
